#!/usr/bin/env python3
"""
Generate a PowerPoint presentation for Transit Tracker project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_transit_tracker_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define colors
    primary_color = RGBColor(59, 130, 246)  # Blue
    secondary_color = RGBColor(16, 185, 129)  # Green
    accent_color = RGBColor(245, 158, 11)  # Orange
    text_color = RGBColor(31, 41, 55)  # Dark gray
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Transit Tracker"
    subtitle.text = "Mobile Driver Dashboard\n\nA comprehensive passenger transportation management system\n\nDeveloped with React.js, TypeScript, and PostgreSQL"
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    content.text = """What is Transit Tracker?

Transit Tracker is a modern web application designed for transit drivers to manage and monitor their routes, passengers, and performance metrics in real-time.

Key Benefits:
• Real-time trip tracking and passenger management
• Automated fare calculation in Kenya Shillings (KES)
• Live analytics and performance monitoring
• Mobile-first responsive design
• Comprehensive revenue tracking"""
    
    # Slide 3: Core Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Core Features"
    content.text = """Essential Functionality

🚌 Active Trip Management
• Start, monitor, and end transit trips
• Real-time GPS location tracking

👥 Passenger Tracking
• Board/alight passenger events
• Passenger pickup at drop-off points

💰 Revenue Management
• Automatic fare calculation
• All financial displays in Kenya Shillings (KES)
• Real-time revenue tracking

📊 Analytics Dashboard
• Performance metrics and statistics
• Trip history and reporting"""
    
    # Slide 4: Technical Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Architecture"
    content.text = """Modern Technology Stack

Frontend:
• React 18 with TypeScript for type safety
• Radix UI components with shadcn/ui styling
• Tailwind CSS for responsive design
• Wouter for client-side routing
• TanStack Query for server state management

Backend:
• Node.js with Express.js server
• PostgreSQL database with Drizzle ORM
• Neon Database for cloud hosting
• WebSocket Server for real-time updates

Build & Development:
• Vite for development and bundling
• ESBuild for server-side compilation"""
    
    # Slide 5: Database Schema
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Database Schema"
    content.text = """Comprehensive Data Model

Core Tables:
• trips - Trip data with status, passenger counts, route information
• passengerEvents - Boarding/alighting events with location data
• locations - Named locations with popularity tracking
• analytics - Daily performance metrics and hourly data
• destinationQueues - Vehicle queue management system
• drivers - Driver information with contact details
• vehicles - Fleet management with capacity tracking

Data Integrity:
• Full type safety with Drizzle ORM
• Zod validation for API requests
• Real-time data synchronization"""
    
    # Slide 6: Real-Time Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Real-Time Features"
    content.text = """Live Communication System

WebSocket Integration:
• Real-time location updates from driver devices
• Live passenger count synchronization
• Trip status broadcasting to connected clients
• Queue position updates and notifications
• Automatic reconnection with exponential backoff

Key Benefits:
• Instant updates across all connected devices
• No page refresh required for data updates
• Seamless user experience
• Reliable connection management"""
    
    # Slide 7: User Interface Design
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "User Interface Design"
    content.text = """Mobile-First Experience

Design Principles:
• Touch-friendly interface with material design patterns
• Bottom navigation for easy thumb navigation
• Responsive design optimized for mobile devices
• Clean, intuitive user interface

Key Components:
• Dashboard with active trip overview
• Trip details with comprehensive information
• Passenger boarding interface
• Analytics dashboard with charts
• Queue management system"""
    
    # Slide 8: Trip Management Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Trip Management Workflow"
    content.text = """Streamlined Operations

Trip Lifecycle:
1. Trip Creation - Select origin, destination, driver, and vehicle
2. Trip Start - Manual trip activation with pending status
3. Passenger Management - Board/alight passengers at any location
4. Route Tracking - Real-time GPS location updates
5. Passenger Pickup - Pick up additional passengers at drop-off points
6. Trip Completion - Manual trip ending with automatic queue entry

Automated Features:
• Automatic fare calculation based on route mappings
• Real-time revenue updates
• Queue position management
• Analytics data generation"""
    
    # Slide 9: Advanced Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Advanced Features"
    content.text = """Enhanced Functionality

Passenger Pickup System:
• Dedicated boarding page instead of popup modals
• Multiple destination groups per pickup
• Individual passenger counts per destination
• Real-time fare calculation
• Automatic revenue tracking

Route Management:
• Predefined route mappings between major locations
• Automatic drop-off point generation
• Dynamic fare calculation based on distance
• Support for multiple pickup/drop-off combinations"""
    
    # Slide 10: Analytics & Reporting
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Analytics & Reporting"
    content.text = """Data-Driven Insights

Key Performance Indicators:
• Total trips and passenger counts
• Revenue tracking in Kenya Shillings
• Distance covered and trip duration
• Driver performance metrics

Visual Analytics:
• Interactive hourly passenger flow charts
• Popular routes analysis with progress bars
• Top performing drivers leaderboard
• Smart insights with peak hours identification

Reporting Features:
• Daily trip summaries
• Real-time analytics refresh (30-60 seconds)
• Historical data analysis
• Performance trend tracking"""
    
    # Slide 11: Fleet Management
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Fleet Management"
    content.text = """Comprehensive Vehicle Tracking

Vehicle Information:
• Number plate, make, model, and year
• Passenger capacity tracking
• Active vehicle status management
• Integration with trip assignments

Driver Management:
• Driver contact information
• Assistant details and contact
• Performance tracking
• Trip assignment history

Queue System:
• Automatic queue entry upon trip completion
• First-come-first-served positioning
• Real-time queue status updates
• Estimated boarding times"""
    
    # Slide 12: Technical Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Implementation"
    content.text = """Development Approach

Code Quality:
• Full TypeScript implementation for type safety
• Component-based architecture with React
• API-first design with RESTful endpoints
• Real-time WebSocket communication

Data Flow:
1. Client connects to WebSocket server
2. Location updates sent from driver device
3. Server broadcasts updates to all connected clients
4. Database updated with location and passenger data
5. Analytics refreshed in real-time

Deployment:
• Replit-optimized development environment
• PostgreSQL database with connection pooling
• Automated build process with Vite and ESBuild"""
    
    # Slide 13: Project Evolution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Evolution"
    content.text = """Development Timeline

Recent Enhancements (July 2025):
✓ Migrated from Replit Agent to standard environment
✓ PostgreSQL database setup with full data persistence
✓ Enhanced trip management with driver/vehicle selection
✓ Passenger pickup functionality at drop-off points
✓ Dedicated passenger boarding pages
✓ Fixed API integration and data validation
✓ Disabled auto-trip completion for better user control
✓ Currency standardization to Kenya Shillings (KES)

Key Achievements:
• Complete database migration and data integrity
• Enhanced user experience with modal-to-page transitions
• Robust error handling and validation
• Real-time communication optimization"""
    
    # Slide 14: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Enhancements"
    content.text = """Roadmap & Opportunities

Potential Improvements:
• GPS integration for automatic location tracking
• Push notifications for queue updates
• Mobile app development (React Native)
• Advanced route optimization algorithms
• Integration with payment systems
• Multi-language support
• Offline functionality for poor connectivity areas

Scalability Features:
• Multi-tenant support for different transport companies
• Advanced reporting and business intelligence
• API integration with external mapping services
• Machine learning for route optimization"""
    
    # Slide 15: Technical Specifications
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Specifications"
    content.text = """System Requirements & Performance

Performance Metrics:
• Sub-second response times for API calls
• Real-time WebSocket updates
• Automatic reconnection handling
• Optimized database queries with connection pooling

Browser Compatibility:
• Modern web browsers (Chrome, Firefox, Safari, Edge)
• Mobile responsive design
• Progressive Web App capabilities

Security Features:
• Type-safe API validation with Zod
• Secure database connections
• Environment variable management
• Input sanitization and validation"""
    
    # Slide 16: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Summary"
    content.text = """Transit Tracker Delivers:

✓ Complete transportation management solution
✓ Real-time passenger and trip tracking
✓ Modern, mobile-first user interface
✓ Comprehensive analytics and reporting
✓ Robust technical architecture
✓ Currency support for Kenya Shillings (KES)

Impact:
• Streamlined driver operations
• Improved passenger experience
• Data-driven decision making
• Scalable and maintainable codebase
• Ready for production deployment

Next Steps:
• User testing and feedback integration
• Performance optimization
• Feature enhancement based on real-world usage"""
    
    # Slide 17: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = """Questions & Discussion

Technology Stack: React + TypeScript + PostgreSQL
Real-time Features: WebSocket Integration
Database: Neon PostgreSQL with Drizzle ORM

Key Features to Demonstrate:
1. Trip creation and management
2. Real-time passenger boarding
3. Analytics dashboard
4. Queue management system
5. Revenue tracking in Kenya Shillings"""
    
    return prs

if __name__ == "__main__":
    try:
        print("Creating Transit Tracker PowerPoint presentation...")
        presentation = create_transit_tracker_presentation()
        
        filename = "Transit_Tracker_Presentation.pptx"
        presentation.save(filename)
        print(f"✓ PowerPoint presentation saved as: {filename}")
        print(f"✓ Total slides: {len(presentation.slides)}")
        print("✓ Presentation ready for download!")
        
    except Exception as e:
        print(f"Error creating presentation: {e}")
        print("Make sure python-pptx is installed: pip install python-pptx")